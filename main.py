# main.py
from fastapi import FastAPI
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from fastapi import HTTPException
from fastapi import WebSocket, WebSocketDisconnect
from typing import Dict
import os
from dotenv import load_dotenv

import asyncio
import json
# Load environment variables
load_dotenv()

# Import your models
from models.requests import StatelessChatRequest, DocumentsEmbedRequest, SummaryRequest,SectionRequest,PlanRequest

# Import your orchestrator
from services.orchestrator import NursingTutor

# Import Firebase initialization
import firebase_admin
from firebase_admin import credentials

# embedding
from langchain_openai import OpenAIEmbeddings
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain.schema import Document

# to store files temporarily
import tempfile
import os
import platform
import requests

# models
from models.reponses import GenerateTitleResponse
from models.requests import GenerateTitleRequest

# document loader
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    CSVLoader,
    Docx2txtLoader,
    UnstructuredExcelLoader,
    UnstructuredPowerPointLoader)

from fastapi import File, UploadFile, Form
from typing import List
from uuid import uuid4
import mimetypes

from core.imageloader import OCRImageLoader
from core.pdfloader import OCRPDFLoader, is_scanned_pdf

# langchain
from langchain.prompts import PromptTemplate
from langchain_openai import ChatOpenAI
from langchain_core.output_parsers import StrOutputParser

# firebase access
import firebase_admin
from firebase_admin import credentials,storage

from core.language import LanguageDetector

from services.vectorstore_manager import vectorstore_manager

import json

# Custom PowerPoint loader that doesn't need NLTK
class SimplePowerPointLoader:
    """Lightweight PowerPoint loader using python-pptx directly."""
    def __init__(self, file_path: str):
        self.file_path = file_path
    
    def load(self):
        """Extract text from PowerPoint slides."""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            raise ImportError("python-pptx is required. Install: pip install python-pptx")
        
        prs = Presentation(self.file_path)
        documents = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            text_content = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                # Regular text shapes
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
                
                # Handle tables - use shape_type to check safely
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    try:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells if cell.text.strip()
                            )
                            if row_text:
                                text_content.append(row_text)
                    except Exception as e:
                        print(f"⚠️ Failed to extract table: {e}")
                        continue
                
                # Handle grouped shapes (recursively extract text)
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    try:
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, "text") and sub_shape.text.strip():
                                text_content.append(sub_shape.text.strip())
                    except Exception as e:
                        print(f"⚠️ Failed to extract grouped shape: {e}")
                        continue
            
            # Create document if there's content
            if text_content:
                page_content = "\n\n".join(text_content)
                doc = Document(
                    page_content=page_content,
                    metadata={
                        "source": os.path.basename(self.file_path),
                        "page": slide_num,
                        "total_slides": len(prs.slides)
                    }
                )
                documents.append(doc)
        
        print(f"✅ Extracted text from {len(documents)} slides")
        return documents

# to load documents
def get_loader_for_file(path):
    ext = os.path.splitext(path)[-1].lower()        
    if ext == ".pdf":# pdf file support
        #Detect if PDF is scanned or text-based
        if is_scanned_pdf(path):
            print("📷 Detected scanned PDF - Using OCR")
            return OCRPDFLoader(path)
        else:
            print("📄 Detected text-based PDF - Using standard loader")
            return PyPDFLoader(path)
    elif ext == ".txt": # txt file support
        return TextLoader(path)
    elif ext == ".csv":  #excel support
        return CSVLoader(path)
    elif ext in [".doc", ".docx"]: # word document support
        return Docx2txtLoader(path) 
    elif ext in [".xls", ".xlsx"]: #excel support
        return UnstructuredExcelLoader(path) 
    elif ext in [".ppt", ".pptx"]: # power point support
        return SimplePowerPointLoader(path)
    elif ext in [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp",".heic"]:  # Add image support
        print("Extracting text from image")
        return OCRImageLoader(path)
    else:
        raise ValueError("Unsupported file type")

# Initialize Firebase
cred = credentials.Certificate("FireBaseAccess.json")
firebase_admin.initialize_app(cred, {
    "storageBucket": os.getenv("FIREBASE_BUCKET", "docai-efb03.firebasestorage.app")
})

# cred = credentials.Certificate("service-account-key.json")  # Path to key file
#     firebase_admin.initialize_app(cred, {
#         'storageBucket': 'docai-efb03.firebasestorage.app'
# })


# Set Google Cloud credentials to the Firebase service account file
# This enables Google Cloud TTS and other Google Cloud APIs
os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = "FireBaseAccess.json"
api_key = os.getenv("OPENAI_API_KEY")
if api_key:
    print(f"✅ OPENAI_API_KEY found")
    print(f"   Length: {len(api_key)}")
    print(f"   First 15 chars: {api_key[:15]}")
    print(f"   Last 4 chars: ...{api_key[-4:]}")
    print(f"   Type: {type(api_key)}")
    # Move the check outside the f-string
    has_whitespace = ' ' in api_key or '\n' in api_key or '\t' in api_key
    print(f"   Contains whitespace: {has_whitespace}")
else:
    print("❌ OPENAI_API_KEY is NOT SET")

# ============================================================================
# COST OPTIMIZATION: Lifespan context manager for startup/shutdown
# ============================================================================
from contextlib import asynccontextmanager

@asynccontextmanager
async def lifespan(_app: FastAPI):
    """Lifespan context manager for startup and shutdown events"""
    # Startup
    asyncio.create_task(periodic_cleanup())
    print("🚀 Started periodic cleanup task (30s interval, 60s timeout)")
    yield
    # Shutdown - cleanup all sessions
    print("🛑 Shutting down - cleaning up all sessions...")
    for chat_id in list(ACTIVE_SESSIONS.keys()):
        cleanup_session(chat_id)

# Create FastAPI app
app = FastAPI(
    title="Nursing Tutor AI",
    version="1.0.0",
    description="AI-powered nursing education assistant with tool calling",
    lifespan=lifespan
)

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "http://localhost:3001",
        "https://docai-efb03.web.app",
        "https://docai-efb03.firebaseapp.com",
        "https://chats.nursequizai.com",
        "https://nursequizai.com",
        "https://ragfastapi-1075876064685.europe-west1.run.app"
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Global session storage
ACTIVE_SESSIONS: Dict[str, NursingTutor] = {}
SESSION_LAST_ACTIVITY: Dict[str, float] = {}  # Track last activity time for each session

# ============================================================================
# COST OPTIMIZATION: Session cleanup configuration
# ============================================================================
SESSION_IDLE_TIMEOUT = 300  # 5 minutes - keep session alive for rapid interactions
CONNECTION_IDLE_TIMEOUT = 300  # 5 minutes - WebSocket idle timeout (not per-message)
SESSION_MAX_AGE = 1800  # 30 minutes - max session lifetime regardless of activity

import time

def cleanup_session(chat_id: str):
    """Clean up session and free memory"""
    if chat_id in ACTIVE_SESSIONS:
        try:
            session = ACTIVE_SESSIONS[chat_id]
            # Clear vectorstore to free memory
            if hasattr(session, 'session') and hasattr(session.session, 'vectorstore'):
                session.session.vectorstore = None
            del ACTIVE_SESSIONS[chat_id]
            print(f"🧹 Cleaned up session for chat_id: {chat_id}")
        except Exception as e:
            print(f"⚠️ Error cleaning up session {chat_id}: {e}")

    if chat_id in SESSION_LAST_ACTIVITY:
        del SESSION_LAST_ACTIVITY[chat_id]

def update_session_activity(chat_id: str):
    """Update last activity timestamp for a session"""
    SESSION_LAST_ACTIVITY[chat_id] = time.time()

# Connection manager for WebSocket connections
class ConnectionManager:
    def __init__(self):
        self.active_connections: Dict[str, WebSocket] = {}
        self.cancellation_flags: Dict[str, bool] = {}  # Track if chat should be cancelled
        self.connection_times: Dict[str, float] = {}  # Track when connection was established
        self.last_activity: Dict[str, float] = {}  # Track last message time

    async def connect(self, websocket: WebSocket, chat_id: str):
        await websocket.accept()
        self.active_connections[chat_id] = websocket
        self.cancellation_flags[chat_id] = False  # Reset cancellation flag
        self.connection_times[chat_id] = time.time()
        self.last_activity[chat_id] = time.time()
        update_session_activity(chat_id)
        print(f"✅ WebSocket connected for chat_id: {chat_id}")

    def disconnect(self, chat_id: str):
        if chat_id in self.active_connections:
            del self.active_connections[chat_id]
            print(f"❌ WebSocket disconnected for chat_id: {chat_id}")
        if chat_id in self.cancellation_flags:
            del self.cancellation_flags[chat_id]
        if chat_id in self.connection_times:
            del self.connection_times[chat_id]
        if chat_id in self.last_activity:
            del self.last_activity[chat_id]

        # COST OPTIMIZATION: DON'T cleanup session immediately on disconnect
        # Keep session alive for 60s so user can reconnect without reloading vectorstore
        # The periodic_cleanup task will clean it up if truly idle
        # This saves cost by avoiding repeated Firebase/vectorstore loads
        update_session_activity(chat_id)  # Reset timer on disconnect

    def update_activity(self, chat_id: str):
        """Update last activity time for a connection"""
        self.last_activity[chat_id] = time.time()
        update_session_activity(chat_id)

    def cancel_stream(self, chat_id: str):
        """Mark a chat's stream for cancellation"""
        self.cancellation_flags[chat_id] = True
        print(f"🛑 Stream cancellation requested for chat_id: {chat_id}")

    def is_cancelled(self, chat_id: str) -> bool:
        """Check if a chat's stream has been cancelled"""
        return self.cancellation_flags.get(chat_id, False)

    def reset_cancellation(self, chat_id: str):
        """Reset cancellation flag for a chat"""
        self.cancellation_flags[chat_id] = False

    async def send_message(self, chat_id: str, message: dict):
        if chat_id in self.active_connections:
            try:
                await self.active_connections[chat_id].send_text(json.dumps(message))
            except Exception as e:
                print(f"Error sending message to {chat_id}: {e}")
                self.disconnect(chat_id)

    async def cleanup_idle_connections(self):
        """Close connections that have been idle too long - called periodically"""
        current_time = time.time()
        to_close = []

        for chat_id, last_time in list(self.last_activity.items()):
            idle_seconds = current_time - last_time
            if idle_seconds > CONNECTION_IDLE_TIMEOUT:
                to_close.append(chat_id)
                print(f"⏰ Connection {chat_id} idle for {idle_seconds:.0f}s, closing...")

        for chat_id in to_close:
            try:
                ws = self.active_connections.get(chat_id)
                if ws:
                    await ws.close(1000, "Idle timeout - reconnect when needed")
            except Exception as e:
                print(f"Error closing idle connection {chat_id}: {e}")
            finally:
                self.disconnect(chat_id)

# Global connection manager
manager = ConnectionManager()

# Set the manager reference in quiztools for cancellation checks
from tools.quiztools import set_connection_manager, get_chat_context_from_db
set_connection_manager(manager)

# ============================================================================
# COST OPTIMIZATION: Background cleanup task
# ============================================================================
async def periodic_cleanup():
    """Background task to clean up idle sessions and connections"""
    while True:
        try:
            await asyncio.sleep(30)  # Check every 30 seconds for 60s timeout
            current_time = time.time()

            # Clean up idle connections
            await manager.cleanup_idle_connections()

            # Clean up expired sessions (even without active WebSocket)
            sessions_to_cleanup = []
            for chat_id, last_activity in list(SESSION_LAST_ACTIVITY.items()):
                idle_time = current_time - last_activity
                if idle_time > SESSION_IDLE_TIMEOUT:
                    sessions_to_cleanup.append(chat_id)
                    print(f"🧹 Session {chat_id} expired (idle {idle_time:.0f}s)")

            for chat_id in sessions_to_cleanup:
                cleanup_session(chat_id)

            # Log stats
            active_sessions = len(ACTIVE_SESSIONS)
            active_connections = len(manager.active_connections)
            if active_sessions > 0 or active_connections > 0:
                print(f"📊 Active: {active_sessions} sessions, {active_connections} connections")

        except Exception as e:
            print(f"⚠️ Cleanup task error: {e}")

# WebSocket endpoint
@app.websocket("/ws/{chat_id}")
async def websocket_endpoint(websocket: WebSocket, chat_id: str):
    #connect
    await manager.connect(websocket, chat_id)
    try:
        while True:
            # Wait for incoming message from client with timeout
            # COST OPTIMIZATION: Use asyncio.wait_for to enforce connection timeout
            try:
                data = await asyncio.wait_for(
                    websocket.receive_text(),
                    timeout=CONNECTION_IDLE_TIMEOUT
                )
            except asyncio.TimeoutError:
                print(f"⏰ WebSocket {chat_id} timed out after {CONNECTION_IDLE_TIMEOUT}s idle")
                await websocket.close(1000, "Idle timeout")
                break

            message = json.loads(data)

            # Update activity timestamp
            manager.update_activity(chat_id)

            # Handle different message types
            await handle_websocket_message(chat_id, message, websocket)

    except WebSocketDisconnect:
        manager.disconnect(chat_id)
        print(f"Client {chat_id} disconnected")
    except Exception as e:
        print(f"WebSocket error for {chat_id}: {e}")
        manager.disconnect(chat_id)
        
async def handle_websocket_message(chat_id: str, message: dict, websocket: WebSocket):
    """Handle incoming WebSocket messages"""
    message_type = message.get("type")
    print(f"📨 Received WebSocket message: type={message_type}, chat_id={chat_id}")

    # check if Im getting a message
    if message_type == "chat_message":
        # Handle regular chat messages (proceed to call the AI Tutor)
        await process_chat_message(chat_id, message, websocket)

    # check if Im getting a ping to get the session alive
    elif message_type == "ping":
        # Handle ping/pong for connection keepalive
        await websocket.send_text(json.dumps({"type": "pong"}))

    # check if user wants to cancel ongoing streaming
    elif message_type == "cancel_stream":
        print(f"🛑 Received cancel request for chat {chat_id}")
        manager.cancel_stream(chat_id)
        await websocket.send_text(json.dumps({
            "type": "stream_cancelled",
            "message": "Stream cancellation requested"
        }))

    # ============================================
    # GAME MODE HANDLERS
    # These handle the gamified quiz flow where users
    # collect serum by answering questions correctly
    # ============================================

    elif message_type == "game_quiz":
        # User wants to start a quiz game
        # This streams questions one at a time
        await process_game_quiz(chat_id, message, websocket)

    elif message_type == "game_deliver":
        # User finished quiz and wants to deliver serum
        # Check if they have enough to save the child
        await process_game_deliver(chat_id, message, websocket)

    elif message_type == "game_retry":
        # User didn't have enough serum, wants to try again
        # Serum persists across retries
        await process_game_retry(chat_id, message, websocket)

    # ============================================
    # MICRO-RATIONALE HANDLER
    # Generate short, encouraging feedback after quiz answer
    # Uses GPT-4.1-nano for speed, with instant fallback
    # ============================================
    elif message_type == "micro_rationale_request":
        await process_micro_rationale(chat_id, message, websocket)

    else:
        # Unknown message type
        await websocket.send_text(json.dumps({
            "type": "error",
            "message": f"Unknown message type: {message_type}"
        }))
               
async def process_chat_message(chat_id: str, message: dict, websocket: WebSocket):
    """Process chat messages through the existing NursingTutor"""
    try:
        # Extract message data
        user_input = message.get("input", "")

        # Get or create session (same as existing logic)
        session_existed = chat_id in ACTIVE_SESSIONS
        print(f"🔍 ACTIVE_SESSIONS keys: {list(ACTIVE_SESSIONS.keys())}")
        print(f"🔍 Looking for chat_id: {chat_id}")
        print(f"🔍 Session exists: {session_existed}")

        if not session_existed:
            ACTIVE_SESSIONS[chat_id] = NursingTutor(chat_id)
            # Still reload insights in case new files were uploaded
            await ACTIVE_SESSIONS[chat_id].load_file_insights_from_firebase()
            print(f"🆕 Created new session for chat {chat_id}")
        else:
            print(f"♻️ Reusing existing session for chat {chat_id}")

        nursing_tutor = ACTIVE_SESSIONS[chat_id]

        # Debug: Check vectorstore status
        has_vectorstore = nursing_tutor.session.vectorstore is not None
        print(f"📊 Session vectorstore status: {'EXISTS in memory' if has_vectorstore else 'NOT in memory'}")
        print(f"📊 Session object id: {id(nursing_tutor.session)}")

        # Ensure vectorstore is loaded (needed for mindmap, quiz, etc.)
        if nursing_tutor.session.vectorstore is None:
            print(f"📥 Loading vectorstore from Firebase for chat {chat_id}...")
            loaded_vectorstore = await vectorstore_manager.load_combined_vectorstore_from_firebase(chat_id)
            if loaded_vectorstore:
                nursing_tutor.session.vectorstore = loaded_vectorstore
                print(f"✅ Vectorstore loaded successfully for {chat_id}")
            else:
                print(f"⚠️ No vectorstore found in Firebase for {chat_id} - background upload may still be in progress")

        # Get chat history for context-aware language detection
        full_context_from_db = await get_chat_context_from_db(chat_id)
        chat_history = full_context_from_db.get("conversation", [])[-10:] if full_context_from_db else []

        # Detect language with chat context
        language = await LanguageDetector.detect_language(user_input, chat_history)
        print(f"Input was entered in language: {language}")
        
        # Send status update
        await websocket.send_text(json.dumps({
            "type": "status",
            "status": "processing",
            "message": "Processing your message..."
        }))
        
        # Process message and stream responses
        async for chunk in nursing_tutor.process_message(user_input, language):
            # Check for cancellation before processing each chunk
            if manager.is_cancelled(chat_id):
                print(f"🛑 Stream cancelled for chat {chat_id}, stopping...")
                await websocket.send_text(json.dumps({
                    "type": "stream_cancelled",
                    "message": "Streaming stopped by user"
                }))
                manager.reset_cancellation(chat_id)  # Reset for next message
                return  # Stop streaming

            # Parse the existing streaming format
            try:
                chunk_data = json.loads(chunk.strip())

                # Debug: Log mindmap-related chunks
                if chunk_data.get("status") in ["mindmap_generating", "mindmap_complete"]:
                    print(f"🧠 Sending mindmap chunk: status={chunk_data.get('status')}, has_data={bool(chunk_data.get('mindmap_data'))}")

                # Forward to WebSocket with type wrapper
                await websocket.send_text(json.dumps({
                    "type": "stream_chunk",
                    "data": chunk_data
                }))

            except json.JSONDecodeError:
                # Handle non-JSON chunks
                await websocket.send_text(json.dumps({
                    "type": "stream_chunk",
                    "data": {"answer_chunk": chunk.strip()}
                }))
        
        # Send completion signal
        manager.reset_cancellation(chat_id)  # Reset cancellation flag
        await websocket.send_text(json.dumps({
            "type": "stream_complete",
            "message": "Response complete"
        }))

        # Update activity timestamp - connection stays open for more messages
        # Will be closed by idle timeout (5 min) or client disconnect
        manager.update_activity(chat_id)
        print(f"✅ Stream complete for {chat_id}, connection stays open for follow-up messages")

    except Exception as e:
        print(f"Error processing chat message: {e}")
        await websocket.send_text(json.dumps({
            "type": "error",
            "message": f"Processing failed: {str(e)}"
        }))


# ============================================================================
# UNUSED FOR NOW, WE GAVE UP ON THE GAMIFICATION IDEA, Code kept in case there is an opportunity for this in the future
# GAME MODE FUNCTIONS
# These functions handle the gamified quiz flow where users collect serum
# by answering NCLEX-style questions correctly to save a sick child.
#
# Flow:
#   1. User uploads file → documents get embedded (existing flow)
#   2. Frontend sends "game_quiz" → backend streams questions
#   3. User answers questions → frontend validates client-side
#   4. Frontend sends "game_deliver" with serum count
#   5. If enough serum → child saved! If not → "game_retry"
#
# Serum Math:
#   - Each correct answer = 20mL
#   - Required to save child = 100mL (default)
#   - Serum ACCUMULATES across retries (forgiving design)
# ============================================================================

async def process_game_quiz(chat_id: str, message: dict, websocket: WebSocket):
    """
    Generate and stream quiz questions for game mode.

    Each question includes the answer so the frontend can validate locally.
    This gives instant feedback without server round-trips.

    Message format from frontend:
    {
        "type": "game_quiz",
        "questionCount": 5,       # Optional, default 5
        "difficulty": "medium",   # Optional, default "medium"
        "existingTopics": []      # Optional, user's existing topics for smart matching
    }
    """
    try:
        # Import what we need
        # Use the Question Bank-integrated version for instant delivery + bank enrichment
        from services.quiz_with_bank import stream_quiz_with_bank as stream_quiz_questions
        from firebase_admin import firestore
        from uuid import uuid4

        db = firestore.client()
        chat_ref = db.collection("chats").document(chat_id)

        print(f"🎮 Starting game quiz for chat: {chat_id}")

        # ------------------------------------------
        # Step 1: Get or create the NursingTutor session
        # This gives us access to the vectorstore with user's documents
        # ------------------------------------------
        if chat_id not in ACTIVE_SESSIONS:
            ACTIVE_SESSIONS[chat_id] = NursingTutor(chat_id)
            await ACTIVE_SESSIONS[chat_id].load_file_insights_from_firebase()

        nursing_tutor = ACTIVE_SESSIONS[chat_id]
        session = nursing_tutor.session  # PersistentSessionContext

        # ------------------------------------------
        # Step 1.5: Ensure vectorstore is loaded from Firebase
        # The upload saves vectorstore to Firebase, but a new WebSocket
        # session needs to load it back into memory.
        #
        # RACE CONDITION HANDLING: The background upload might still be
        # in progress when this runs, so we retry a few times with delays.
        # ------------------------------------------
        if session.vectorstore is None:
            print(f"📥 Loading vectorstore from Firebase for game quiz...")

            # Notify frontend we're loading documents
            await websocket.send_text(json.dumps({
                "type": "stream_chunk",
                "data": {
                    "status": "game_loading_documents",
                    "message": "Loading your documents..."
                }
            }))

            # Retry up to 5 times with increasing delays (total ~10 seconds max)
            max_retries = 5
            retry_delays = [1, 2, 2, 3, 3]  # seconds between retries
            loaded_vectorstore = None

            for attempt in range(max_retries):
                loaded_vectorstore = await vectorstore_manager.load_combined_vectorstore_from_firebase(chat_id)

                if loaded_vectorstore:
                    session.vectorstore = loaded_vectorstore
                    print(f"✅ Vectorstore loaded successfully for {chat_id} (attempt {attempt + 1})")
                    break
                else:
                    if attempt < max_retries - 1:
                        wait_time = retry_delays[attempt]
                        print(f"⏳ Vectorstore not ready yet, waiting {wait_time}s... (attempt {attempt + 1}/{max_retries})")
                        await asyncio.sleep(wait_time)
                    else:
                        print(f"⚠️ No vectorstore found in Firebase for {chat_id} after {max_retries} attempts")

            if not loaded_vectorstore:
                await websocket.send_text(json.dumps({
                    "type": "error",
                    "message": "Documents are still processing. Please wait a moment and try again."
                }))
                return

        # ------------------------------------------
        # Step 2: Get or initialize game state from Firestore
        # ------------------------------------------
        chat_doc = chat_ref.get()

        if not chat_doc.exists:
            # This shouldn't happen - chat should exist from file upload
            await websocket.send_text(json.dumps({
                "type": "error",
                "message": "Chat not found. Please upload a file first."
            }))
            return

        chat_data = chat_doc.to_dict()
        game_state = chat_data.get("gameState", {})

        # If no game state exists, initialize it
        if not game_state:
            game_state = {
                "status": "in_progress",
                "serumCollected": 0,
                "serumRequired": 100,
                "attempts": 0,
                "completedAt": None
            }
            chat_ref.update({"gameState": game_state})
            print(f"🎮 Initialized new game state for chat: {chat_id}")

        # ------------------------------------------
        # Step 3: Send game initialized message
        # Frontend uses this to set up the UI
        # ------------------------------------------
        await websocket.send_text(json.dumps({
            "type": "stream_chunk",
            "data": {
                "status": "game_initialized",
                "serumCollected": game_state.get("serumCollected", 0),
                "serumRequired": game_state.get("serumRequired", 100)
            }
        }))

        # ------------------------------------------
        # Step 4: Extract quiz parameters from message
        # ------------------------------------------
        question_count = message.get("questionCount", 5)
        difficulty = message.get("difficulty", "medium")
        existing_topics = message.get("existingTopics", [])  # User's existing topics
        quiz_id = f"quiz_{uuid4().hex[:8]}"

        print(f"🎮 Generating {question_count} {difficulty} questions...")
        print(f"📚 User's existing topics: {existing_topics}")
        print(f"📊 Session vectorstore status: {session.vectorstore is not None}")
        if session.vectorstore:
            print(f"📊 Vectorstore type: {type(session.vectorstore)}")

        # ------------------------------------------
        # Step 5: Stream questions using existing function
        # Each question includes answer + justification for client-side validation
        # ------------------------------------------
        question_index = 0

        async for chunk in stream_quiz_questions(
            topic="",                    # Empty = use all document content
            difficulty=difficulty,
            num_questions=question_count,
            source="documents",          # Generate from user's uploaded docs
            session=session,
            empathetic_message=None,     # No intro message for game mode
            chat_id=chat_id,             # For cancellation checking
            existing_topics=existing_topics  # User's existing topics for smart matching
        ):
            # Handle different chunk types from stream_quiz_questions

            if chunk.get("status") == "generating":
                # Progress update: "Generating question 2 of 5..."
                await websocket.send_text(json.dumps({
                    "type": "stream_chunk",
                    "data": {
                        "status": "game_generating",
                        "current": chunk.get("current"),
                        "total": chunk.get("total")
                    }
                }))

            elif chunk.get("status") == "question_ready":
                # A complete question is ready to send
                question = chunk.get("question")

                # Send the full question (including answer for client-side validation)
                await websocket.send_text(json.dumps({
                    "type": "stream_chunk",
                    "data": {
                        "status": "game_question_ready",
                        "question": {
                            "index": question_index,
                            "question": question.get("question"),
                            "options": question.get("options"),
                            "answer": question.get("answer"),           # For client validation
                            "justification": question.get("justification"),  # Show after answer
                            "topic": question.get("topic", "General"),
                            "serumValue": 20  # Each correct answer = 20mL
                        },
                        "quizId": quiz_id,
                        "isFirst": question_index == 0  # Frontend can start showing UI
                    }
                }))

                question_index += 1
                print(f"✅ Sent question {question_index}/{question_count}")

            elif chunk.get("status") == "quiz_complete":
                # All questions generated
                await websocket.send_text(json.dumps({
                    "type": "stream_chunk",
                    "data": {
                        "status": "game_quiz_complete",
                        "totalQuestions": chunk.get("total_generated", question_count)
                    }
                }))
                print(f"🎮 Quiz complete! Sent {question_index} questions")

                # Update activity - connection stays open for delivery/retry
                manager.update_activity(chat_id)
                print(f"✅ Game quiz complete for {chat_id}, connection stays open")

    except Exception as e:
        print(f"❌ Game quiz error: {e}")
        import traceback
        traceback.print_exc()
        await websocket.send_text(json.dumps({
            "type": "error",
            "message": f"Failed to generate quiz: {str(e)}"
        }))


async def process_game_deliver(chat_id: str, message: dict, websocket: WebSocket):
    """
    Handle serum delivery attempt.

    Frontend reports how much serum was collected (based on correct answers).
    We check if it's enough to save the child.

    Message format from frontend:
    {
        "type": "game_deliver",
        "serumCollected": 80  # Total serum from this session
    }
    """
    try:
        from firebase_admin import firestore
        db = firestore.client()

        chat_ref = db.collection("chats").document(chat_id)
        chat_doc = chat_ref.get()

        if not chat_doc.exists:
            await websocket.send_text(json.dumps({
                "type": "error",
                "message": "Game not found"
            }))
            return

        # ------------------------------------------
        # Step 1: Get current game state
        # ------------------------------------------
        chat_data = chat_doc.to_dict()
        game_state = chat_data.get("gameState", {})

        # Get the serum values
        serum_from_this_quiz = message.get("serumCollected", 0)
        previous_serum = game_state.get("serumCollected", 0)
        serum_required = game_state.get("serumRequired", 100)

        # Total serum = previous attempts + this attempt
        total_serum = previous_serum + serum_from_this_quiz

        print(f"🧪 Delivery attempt: {serum_from_this_quiz}mL this quiz + {previous_serum}mL previous = {total_serum}mL total")
        print(f"🧪 Required: {serum_required}mL")

        # ------------------------------------------
        # Step 2: Update game state with new serum total
        # ------------------------------------------
        attempts = game_state.get("attempts", 0) + 1

        chat_ref.update({
            "gameState.serumCollected": total_serum,
            "gameState.attempts": attempts
        })

        # ------------------------------------------
        # Step 3: Check if we have enough serum
        # ------------------------------------------
        if total_serum >= serum_required:
            # SUCCESS! Child is saved!
            chat_ref.update({
                "gameState.status": "completed",
                "gameState.completedAt": firestore.SERVER_TIMESTAMP
            })

            await websocket.send_text(json.dumps({
                "type": "stream_chunk",
                "data": {
                    "status": "game_child_saved",
                    "serumDelivered": total_serum,
                    "attempts": attempts,
                    "message": "The serum worked. The child is stabilizing. You saved them!"
                }
            }))

            print(f"🎉 Child saved! Total serum: {total_serum}mL in {attempts} attempt(s)")

        else:
            # NOT ENOUGH - Need to retry
            serum_needed = serum_required - total_serum

            await websocket.send_text(json.dumps({
                "type": "stream_chunk",
                "data": {
                    "status": "game_need_more_serum",
                    "serumCollected": total_serum,
                    "serumRequired": serum_required,
                    "serumNeeded": serum_needed,
                    "attempts": attempts,
                    "message": f"The child needs {serum_needed}mL more serum. Keep going!"
                }
            }))

            print(f"⚠️ Need more serum: {serum_needed}mL more required")

    except Exception as e:
        print(f"❌ Game deliver error: {e}")
        import traceback
        traceback.print_exc()
        await websocket.send_text(json.dumps({
            "type": "error",
            "message": f"Delivery failed: {str(e)}"
        }))


async def process_game_retry(chat_id: str, message: dict, websocket: WebSocket):
    """
    Handle retry request after not having enough serum.

    Key behavior: Serum PERSISTS across retries!
    This is intentional - we want to encourage persistence, not punish failure.

    Message format from frontend:
    {
        "type": "game_retry",
        "questionCount": 5,      # Optional
        "difficulty": "medium"   # Optional
    }
    """
    try:
        from firebase_admin import firestore
        db = firestore.client()

        chat_ref = db.collection("chats").document(chat_id)
        chat_doc = chat_ref.get()

        if not chat_doc.exists:
            await websocket.send_text(json.dumps({
                "type": "error",
                "message": "Game not found"
            }))
            return

        # ------------------------------------------
        # Get current serum level (this persists!)
        # ------------------------------------------
        chat_data = chat_doc.to_dict()
        game_state = chat_data.get("gameState", {})
        current_serum = game_state.get("serumCollected", 0)
        serum_required = game_state.get("serumRequired", 100)
        serum_needed = serum_required - current_serum

        print(f"🔄 Retry requested. Current serum: {current_serum}mL, need {serum_needed}mL more")

        # ------------------------------------------
        # Notify frontend that retry is starting
        # ------------------------------------------
        await websocket.send_text(json.dumps({
            "type": "stream_chunk",
            "data": {
                "status": "game_retry_starting",
                "serumCollected": current_serum,
                "serumRequired": serum_required,
                "serumNeeded": serum_needed,
                "message": "Generating more questions... Your serum is safe!"
            }
        }))

        # ------------------------------------------
        # Reuse the game_quiz handler to generate new questions
        # ------------------------------------------
        await process_game_quiz(chat_id, message, websocket)

    except Exception as e:
        print(f"❌ Game retry error: {e}")
        import traceback
        traceback.print_exc()
        await websocket.send_text(json.dumps({
            "type": "error",
            "message": f"Retry failed: {str(e)}"
        }))


# ============================================================================
# MICRO-RATIONALE HANDLER
# Generate short, encouraging feedback using GPT-4.1-nano
# ============================================================================

# Static fallbacks for instant response when GPT-nano is slow/fails
MICRO_RATIONALE_FALLBACKS = {
    "correct": {
        "en": {
            "encouragement": "Nice — that's a core clinical priority.",
            "rationale": "The priority is addressing life-threatening problems first."
        },
        "fr": {
            "encouragement": "Bien joué — c'est une priorité clinique essentielle.",
            "rationale": "La priorité est de traiter d'abord les problèmes vitaux."
        }
    },
    "incorrect": {
        "en": {
            "encouragement": "Good attempt — this is a common exam trap.",
            "rationale": "Focus on the most critical intervention first."
        },
        "fr": {
            "encouragement": "Bonne tentative — c'est un piège d'examen courant.",
            "rationale": "Concentrez-vous d'abord sur l'intervention la plus critique."
        }
    }
}


async def process_micro_rationale(chat_id: str, message: dict, websocket: WebSocket):
    """
    Generate short, encouraging feedback for quiz answers using GPT-4.1-nano.

    Hard constraints:
    - Encouragement: ≤ 12 words
    - Rationale: ≤ 2 sentences
    - No lists, no option-by-option breakdown
    - Must return in <300ms or use fallback

    Message format from frontend:
    {
        "type": "micro_rationale_request",
        "question": "What is the primary purpose...",
        "correct_answer": "To quickly determine...",
        "selected_answer": "To assess the patient's...",
        "is_correct": false,
        "topic": "Emergency Evaluation",
        "original_rationale": "Option C is correct because...",
        "language": "en"
    }
    """
    try:
        # Extract data from message
        question = message.get("question", "")
        correct_answer = message.get("correct_answer", "")
        selected_answer = message.get("selected_answer", "")
        is_correct = message.get("is_correct", False)
        topic = message.get("topic", "General")
        original_rationale = message.get("original_rationale", "")
        language = message.get("language", "en")

        # Ensure language is supported
        lang = language if language in ["en", "fr"] else "en"

        # Get fallback based on correctness
        fallback_key = "correct" if is_correct else "incorrect"
        fallback = MICRO_RATIONALE_FALLBACKS[fallback_key][lang]

        print(f"🎯 Micro-rationale request: {'correct' if is_correct else 'incorrect'} answer")

        try:
            # Use GPT-4.1-nano for fastest response
            llm = ChatOpenAI(model="gpt-4.1-nano", temperature=0.7, request_timeout=0.5)

            # Build tone instruction
            if is_correct:
                tone = "Supportive and confirming. Reinforce why this was the right choice."
            else:
                tone = "Supportive and non-judgmental. Normalize the mistake, explain briefly."

            # Build prompt - minimal, focused
            prompt = f"""Quiz feedback generator. Keep responses SHORT and ENCOURAGING.

Question: {question[:300]}
Correct answer: {correct_answer[:150]}
User selected: {selected_answer[:150]}
User was: {"CORRECT" if is_correct else "INCORRECT"}
Topic: {topic}

{f"Original rationale to condense: {original_rationale[:400]}" if original_rationale else ""}

HARD CONSTRAINTS:
- Encouragement: ≤ 12 words, {tone}
- Rationale: ≤ 2 sentences explaining the key concept
- NO lists, NO "Option A/B/C", NO "According to..."
- Be warm, brief, human

Return ONLY valid JSON:
{{"encouragement": "...", "rationale": "..."}}"""

            # Race against timeout
            result = await asyncio.wait_for(
                llm.ainvoke(prompt),
                timeout=0.5  # 500ms hard limit
            )

            # Parse response
            response_text = result.content.strip()

            # Extract JSON from response
            import re
            if response_text.startswith("{"):
                parsed = json.loads(response_text)
            else:
                # Try to extract JSON from response
                json_match = re.search(r'\{[^}]+\}', response_text)
                if json_match:
                    parsed = json.loads(json_match.group())
                else:
                    raise ValueError("No JSON found in response")

            # Send successful response
            await websocket.send_text(json.dumps({
                "type": "micro_rationale_response",
                "data": {
                    "encouragement": parsed.get("encouragement", fallback["encouragement"])[:100],
                    "rationale": parsed.get("rationale", fallback["rationale"])[:250],
                    "source": "gpt-nano"
                }
            }))

            print(f"✅ Micro-rationale generated via GPT-nano")
            return

        except asyncio.TimeoutError:
            print("⚠️ GPT-nano timeout, using fallback")
        except Exception as e:
            print(f"⚠️ GPT-nano error: {e}, using fallback")

        # Fallback response
        await websocket.send_text(json.dumps({
            "type": "micro_rationale_response",
            "data": {
                **fallback,
                "source": "fallback"
            }
        }))

    except Exception as e:
        print(f"❌ Micro-rationale error: {e}")
        # Return generic fallback on total failure
        await websocket.send_text(json.dumps({
            "type": "micro_rationale_response",
            "data": {
                "encouragement": "Keep going — you're learning!",
                "rationale": "Each question helps build your understanding.",
                "source": "error_fallback"
            }
        }))


# ============================================================================
# SUPPORTING ENDPOINTS (keep your existing ones)
# ============================================================================
def get_temp_dir():
    """Get appropriate temp directory for the current environment"""
    if platform.system() == "Linux":
        # On Cloud Run (Linux), use /tmp which is writable
        return "/tmp"
    else:
        # On Windows/Mac, use system default
        return tempfile.gettempdir()


# ============================================================================
# POST-UPLOAD MESSAGE HELPERS
# ============================================================================
# These functions build the friendly message shown after file upload.
# They use templates instead of LLM calls for instant response + zero cost.
# ============================================================================

def build_post_upload_message(topics: list, file_count: int, language: str) -> str:
    """
    Build a friendly message after file upload using templates.

    WHY TEMPLATES INSTEAD OF LLM:
    - Instant response (no API latency)
    - Zero cost (no tokens used)
    - Predictable output
    - Still feels natural with variety built in

    Args:
        topics: List of topics extracted from uploaded files
        file_count: Number of files uploaded
        language: User's language ('english', 'french', 'fr', etc.)

    Returns:
        A friendly message string
    """
    import random

    # Determine if French
    is_french = language.lower() in ["fr", "french", "français"]

    # Format topics into readable string
    if topics:
        if len(topics) == 1:
            topics_str = topics[0]
        elif len(topics) == 2:
            topics_str = f"{topics[0]} and {topics[1]}" if not is_french else f"{topics[0]} et {topics[1]}"
        else:
            # "Topic1, Topic2, and Topic3"
            if is_french:
                topics_str = ", ".join(topics[:-1]) + f" et {topics[-1]}"
            else:
                topics_str = ", ".join(topics[:-1]) + f", and {topics[-1]}"
    else:
        topics_str = "your study material" if not is_french else "ton matériel d'étude"

    # File count text
    if file_count == 1:
        file_text = "your notes" if not is_french else "tes notes"
    else:
        file_text = f"all {file_count} files" if not is_french else f"les {file_count} fichiers"

    # Message templates - variety keeps it feeling natural
    if is_french:
        templates = [
            f"C'est bon! J'ai parcouru {file_text} et trouvé du contenu sur {topics_str}. Par où veux-tu commencer?",
            f"Parfait! J'ai analysé {file_text} — on a de la matière sur {topics_str}. Qu'est-ce qu'on attaque en premier?",
            f"J'ai tout reçu! Tes documents couvrent {topics_str}. Comment veux-tu étudier?",
        ]
    else:
        templates = [
            f"Got it! I've gone through {file_text} and found material on {topics_str}. Where would you like to start?",
            f"All set! I've looked through {file_text} — there's good content on {topics_str}. How do you want to dive in?",
            f"Nice! Your documents cover {topics_str}. Pick a way to start studying!",
        ]

    return random.choice(templates)


def get_post_upload_actions(language: str) -> list:
    """
    Get the action buttons to show after file upload.

    These are the primary actions a student would want after uploading notes:
    1. Quiz - Test their knowledge
    2. Flashcards - Memorize key concepts
    3. Study sheet - Get a summary breakdown
    4. Audio - Listen to a lecture on the topics

    Args:
        language: User's language

    Returns:
        List of action dictionaries with id, label, and icon
    """
    is_french = language.lower() in ["fr", "french", "français"]

    if is_french:
        return [
            {"id": "quiz", "label": "Teste-moi sur ces sujets", "icon": "🧪"},
            {"id": "flashcards", "label": "Créer des flashcards", "icon": "📇"},
            {"id": "studysheet", "label": "Fais-moi un résumé", "icon": "📝"},
            {"id": "audio", "label": "Écouter une leçon audio", "icon": "🎧"},
            {"id": "mindmap", "label": "Créer une carte mentale", "icon": "🧠"}
        ]
    else:
        return [
            {"id": "quiz", "label": "Quiz me on these topics", "icon": "🧪"},
            {"id": "flashcards", "label": "Create flashcards to study", "icon": "📇"},
            {"id": "studysheet", "label": "Break it down for me", "icon": "📝"},
            {"id": "audio", "label": "Listen to an audio lesson", "icon": "🎧"},
            {"id": "mindmap", "label": "Create a mind map", "icon": "🧠"}
        ]


@app.post("/chat/upload-files")
async def upload_multiple_files(
    files: List[UploadFile] = File(...),
    chat_id: str = Form(...),
    user_id: str = Form(...),
    language: str = Form(...),
):
    """Upload multiple files and process in parallel."""
    # Validate and normalize language parameter
    if not language or language == "undefined" or language == "null":
        language = "english"  # Default fallback

    # Normalize language code (handle 'fr-FR', 'fr-CA', etc.)
    language = language.lower().split('-')[0]  # 'fr-FR' -> 'fr'
    prompt_language = _language_for_prompt(language)
    print(f"*Upload received: user {user_id}, chat:{chat_id}, language:{language}*")
    # Read all files immediately
    file_data_list = []
    for file in files:
        try:
            file_bytes = await file.read()
            file_data_list.append({
                "filename": file.filename,
                "content_type": file.content_type,
                "bytes": file_bytes,
                "size": len(file_bytes)
            })
            print(f"✅ Read file: {file.filename} ({len(file_bytes)} bytes)")
        except Exception as read_error:
            print(f"❌ Failed to read {file.filename}: {read_error}")
            file_data_list.append({
                "filename": file.filename,
                "error": str(read_error)
            })
    
    async def process_and_stream():
        try:
            if not file_data_list:
                yield json.dumps({
                    "type": "error",
                    "message": "No files provided"
                }) + "\n"
                return
            
            # Check for read errors
            failed_reads = [f for f in file_data_list if "error" in f]
            if failed_reads:
                for failed in failed_reads:
                    yield json.dumps({
                        "type": "file_error",
                        "filename": failed["filename"],
                        "message": f"Failed to read file: {failed['error']}"
                    }) + "\n"
            
            # Get successfully read files
            valid_files = [f for f in file_data_list if "bytes" in f]
            
            if not valid_files:
                yield json.dumps({
                    "type": "error",
                    "message": "No valid files to process"
                }) + "\n"
                return
            
            yield json.dumps({
                "type": "batch_start",
                "total_files": len(valid_files),
                "filenames": [f["filename"] for f in valid_files]
            }) + "\n"
            
            # ========================================
            # LOAD EXISTING VECTORSTORE IF NEEDED
            # ========================================
            yield json.dumps({
                "type": "loading_existing_documents",
                "message": "Reading existing documents..."
            }) + "\n"
            
            await ensure_session_with_vectorstore(chat_id)
            
            # ========================================
            # PROCESS FILES (EMBEDDING ONLY)
            # ========================================
            semaphore = asyncio.Semaphore(15)
            
            async def process_single_file_data(file_data):
                async with semaphore:
                    return await process_file_from_bytes(
                        file_data["bytes"],
                        file_data["filename"],
                        chat_id,
                        user_id,
                        language  # Pass browser language
                    )
            
            # Start all file processing tasks
            tasks = [process_single_file_data(fd) for fd in valid_files]
            
            # Stream results as they complete
            total_words = 0
            completed_files = []
            file_documents = {}  # filename -> documents
            file_bytes_map = {}  # filename -> bytes (for background upload)
            
            for coro in asyncio.as_completed(tasks):
                try:
                    result = await coro
                    
                    # Extract updates, documents, insights, and file bytes
                    updates = result.get("updates", [])
                    documents = result.get("documents", [])
                    insights = result.get("insights")  # ← NEW
                    filename = result.get("filename", "unknown")
                    file_bytes = result.get("file_bytes")
                    
                    print(f"🔍 DEBUG: insights type = {type(insights)}")  # ← ADD THIS
                    print(f"🔍 DEBUG: insights value = {insights}") 
                    
                    print(f"🔍 Processing result for: {filename}")
                    print(f"   Documents count: {len(documents)}")
                    
                    if insights:
                        print(f"   Insights extracted from upload: {len(insights.get('topics', []))} topics")
                    
                    # Store documents for background upload
                    if documents:
                        file_documents[filename] = documents
                    if file_bytes:
                        file_bytes_map[filename] = file_bytes
                    
                    # ========================================
                    # 🆕 STORE INSIGHTS IN SESSION
                    # ========================================
                    if insights and chat_id in ACTIVE_SESSIONS:
                        session = ACTIVE_SESSIONS[chat_id]
                        #update the session language using the front-end browser language when uploading
                        session.session.user_language = language 
                        if not hasattr(session.session, "file_insights"):
                            session.session.file_insights = {}
                        session.session.file_insights[filename] = insights
                    
                    # Stream JSON updates to frontend
                    for update in updates:
                        if not isinstance(update, dict):
                            continue
                        
                        yield json.dumps(update) + "\n"
                        
                        if update.get("type") == "file_complete":
                            total_words += update.get("word_count", 0)
                            completed_files.append(update.get("file_id"))
                
                except Exception as e:
                    print(f"❌ Task error: {e}")
                    import traceback
                    traceback.print_exc()
                    yield json.dumps({
                        "type": "file_error",
                        "message": str(e)
                    }) + "\n"
                
            
            # ========================================
            # Generate Upload Summary (1-2 sentences)
            # ========================================
            if chat_id in ACTIVE_SESSIONS:
                session = ACTIVE_SESSIONS[chat_id]
                file_insights = getattr(session.session, "file_insights", {})
                
                if file_insights:
                    try:
                        # Aggregate all insights
                        all_topics = []
                        all_doc_types = []
                        
                        for filename, insights in file_insights.items():
                            if insights:
                                all_topics.extend(insights.get("topics", []))
                                doc_type = insights.get("document_type", "")
                                if doc_type:
                                    all_doc_types.append(doc_type)
                        
                        # Deduplicate
                        unique_topics = list(set(all_topics))
                        unique_doc_types = list(set(all_doc_types))
                        
                        # Generate summary with LLM
                        llm = ChatOpenAI(model="gpt-4o-mini", temperature=0.5)
                        
                        
                        summary_prompt = f"""Generate a brief 1-2 sentence summary about what these uploaded documents contain.

                        Files: {len(file_insights)} document(s)
                        Topics found: {', '.join(unique_topics)}
                        Document types: {', '.join(unique_doc_types) if unique_doc_types else 'various'}

                        Write a natural, conversational summary that tells the student what content was found.
                        Examples:
                        - "I found materials about cardiac pharmacology and arrhythmia management."
                        - "Les documents sur la pharmacologie cardiaque et la gestion des arythmies."

                        IMPORTANT REQUIREMENT: The entire summary must be written in {prompt_language}.

                        Return ONLY the summary text in {prompt_language}, nothing else."""
                        
                        summary_response = await llm.ainvoke([
                            {"role": "user", "content": summary_prompt}
                        ])
                        
                        upload_summary = summary_response.content.strip()
                        
                        # Yield summary to frontend
                        yield json.dumps({
                            "type": "upload_summary",
                            "summary": upload_summary,
                            "file_count": len(valid_files),
                            "filenames": [f["filename"] for f in valid_files]
                        }) + "\n"
                        
                        print(f"📝 Generated upload summary: {upload_summary}")
                        
                        
                        title_prompt = f"""Generate a short, descriptive chat title in 3 to 6 words based on these uploaded study materials.

                       
                        Topics Found: {', '.join(unique_topics) if unique_topics else 'medical content'}
                        Document Types: {', '.join(unique_doc_types) if unique_doc_types else 'study materials'}

                        Requirements:
                        - Be concise and clear
                        - Focus on the main topic/subject area
                        - Max 6 words
                        - Make it specific to the content (e.g., "Cardiac Pharmacology Notes", "NCLEX Respiratory Review")
                        - Write in {prompt_language}

                        Return ONLY the title, no quotes or extra text."""
                        
                        title_response = await llm.ainvoke([
                            {"role": "user", "content": title_prompt}
                        ])
                        
                        chat_title = title_response.content.strip().replace('"', '').replace("'", "")

                        # Update Firebase chat document with new title
                        from firebase_admin import firestore
                        db = firestore.client()

                        db.collection('chats').document(chat_id).update({
                            'title': chat_title,
                            'updatedAt': firestore.SERVER_TIMESTAMP
                        })

                        print(f"✅ Auto-generated and saved chat title: '{chat_title}'")
                        
                    except Exception as summary_error:
                        print(f"⚠️ Summary generation failed: {summary_error}")
            
            # ========================================
            # DISABLED: Suggestions after upload
            # PostUploadActions now handles guiding the user
            # ========================================
            # Suggestions are only generated during normal conversation flow
            # ========================================
            # READY TO CHAT - USER CAN START IMMEDIATELY
            # ========================================
            yield json.dumps({
                "type": "ready_to_chat",
                "message": "Files processed! You can start asking questions.",
                "total_files": len(valid_files),
                "completed_files": len(completed_files),
                "total_words": total_words
            }) + "\n"
            
            # ========================================
            # BACKGROUND: UPLOAD EVERYTHING
            # ========================================
            if chat_id in ACTIVE_SESSIONS:
                session = ACTIVE_SESSIONS[chat_id]
                has_vs = session.session.vectorstore is not None
                print(f"📤 Upload complete for {chat_id}: session in ACTIVE_SESSIONS=True, vectorstore={'EXISTS' if has_vs else 'None'}")
                if session.session.vectorstore:
                    # Start background task (fire-and-forget with retry)
                    asyncio.create_task(
                        upload_everything_background(
                            chat_id=chat_id,
                            vectorstore=session.session.vectorstore,
                            file_documents=file_documents,
                            file_bytes_map=file_bytes_map
                        )
                    )
            else:
                print(f"⚠️ Upload complete but chat_id {chat_id} NOT in ACTIVE_SESSIONS!")
            
            # ========================================
            # FINAL SUMMARY (IMMEDIATE)
            # ========================================
            yield json.dumps({
                "type": "all_complete",
                "total_files": len(valid_files),
                "completed_files": len(completed_files),
                "total_words": total_words,
                "status": "success"
            }) + "\n"

            # ========================================
            # POST-UPLOAD: FRIENDLY MESSAGE + ACTIONS
            # ========================================
            #
            # WHY: After uploading, users see stats but don't know what to do next.
            # This sends a friendly AI message with action buttons to guide them.
            #
            # COST OPTIMIZATION: We reuse the topics already extracted during upload
            # (stored in file_insights) instead of making another LLM call to analyze.
            # The friendly message is built from a template - no extra API call needed.
            #
            # FLOW:
            # 1. Collect topics from all uploaded files (already extracted)
            # 2. Build a friendly message using templates (instant, no LLM)
            # 3. Send to frontend + save to Firebase
            # ========================================

            if chat_id in ACTIVE_SESSIONS:
                session = ACTIVE_SESSIONS[chat_id]
                file_insights = getattr(session.session, "file_insights", {})

                if file_insights:
                    try:
                        # -----------------------------------------
                        # STEP 1: Collect all topics and educational insights from uploaded files
                        # These were already extracted during file processing
                        # -----------------------------------------
                        all_topics = []
                        all_insights = []
                        for filename, insights in file_insights.items():
                            if insights and insights.get("topics"):
                                all_topics.extend(insights.get("topics", []))
                            if insights and insights.get("insights"):
                                all_insights.extend(insights.get("insights", []))

                        # Remove duplicates, keep max 5 topics and 3 insights for readability
                        unique_topics = list(set(all_topics))[:5]
                        # Deduplicate insights by topic
                        seen_topics = set()
                        unique_insights = []
                        for insight in all_insights:
                            topic = insight.get("topic", "")
                            if topic not in seen_topics:
                                seen_topics.add(topic)
                                unique_insights.append(insight)
                            if len(unique_insights) >= 3:
                                break

                        filenames = [f["filename"] for f in valid_files]
                        file_count = len(valid_files)

                        # -----------------------------------------
                        # STEP 2: Build friendly message from template
                        # No LLM call = instant response, zero cost
                        # -----------------------------------------
                        friendly_message = build_post_upload_message(
                            topics=unique_topics,
                            file_count=file_count,
                            language=language
                        )

                        # -----------------------------------------
                        # STEP 3: Build action buttons (localized)
                        # These appear below the message for quick actions
                        # -----------------------------------------
                        actions = get_post_upload_actions(language)

                        # -----------------------------------------
                        # STEP 4: Send to frontend
                        # Frontend will handle saving to Firebase to control timing
                        # (ensures LoadingMessageBox is saved first, then PostUploadActions)
                        # -----------------------------------------
                        yield json.dumps({
                            "type": "post_upload_message",
                            "message": friendly_message,
                            "topics": unique_topics,
                            "insights": unique_insights,  # Educational insights for orientation flow
                            "filenames": filenames,
                            "file_count": file_count,
                            "actions": actions
                        }) + "\n"

                        print(f"✅ Post-upload message sent to frontend (frontend saves to Firebase)")
                        print(f"   Topics: {unique_topics}")
                        print(f"   Educational insights: {len(unique_insights)}")

                    except Exception as post_upload_error:
                        # Non-critical - user can still chat even if this fails
                        print(f"⚠️ Post-upload message failed: {post_upload_error}")

        except Exception as e:
            print(f"❌ Batch processing error: {e}")
            import traceback
            traceback.print_exc()
            yield json.dumps({
                "type": "error",
                "message": str(e)
            }) + "\n"
    
    return StreamingResponse(
        process_and_stream(),
        media_type="application/x-ndjson"
    )

# ============================================================================
# HELPERS FOR FILE UPLOAD START
# ============================================================================
async def upload_everything_background(
    chat_id: str,
    vectorstore: FAISS,
    file_documents: Dict[str, List[Document]],
    file_bytes_map: Dict[str, bytes],
    max_retries: int = 3
):
    """
    Upload files and vectorstores to Firebase in the background.
    Includes automatic retry on failure (silent).
    """
    retry_count = 0
    
    while retry_count < max_retries:
        try:
            print(f"🔄 Background upload starting for chat {chat_id} (attempt {retry_count + 1}/{max_retries})...")
            
            # ========================================
            # UPLOAD FILES TO FIREBASE STORAGE
            # ========================================
            file_upload_tasks = []
            
            for filename, file_bytes in file_bytes_map.items():
                file_upload_tasks.append(
                    firebase_upload_task_simple(file_bytes, filename, chat_id)
                )
            
            # ========================================
            # UPLOAD VECTORSTORES
            # ========================================
            vectorstore_results = await vectorstore_manager.upload_all_vectorstores(
                chat_id=chat_id,
                combined_vectorstore=vectorstore,
                file_documents=file_documents
            )
            
            # ========================================
            # UPLOAD FILES IN PARALLEL
            # ========================================
            file_results = await asyncio.gather(*file_upload_tasks, return_exceptions=True)
            
            # Check results
            file_failures = [r for r in file_results if isinstance(r, Exception)]
            
            if vectorstore_results["combined_success"] and len(file_failures) == 0:
                print(f"✅ Background upload complete for chat {chat_id}")
                return  # Success - exit
            else:
                print(f"⚠️ Background upload had issues (attempt {retry_count + 1})")
                if not vectorstore_results["combined_success"]:
                    print(f"   - Combined vectorstore failed")
                if file_failures:
                    print(f"   - {len(file_failures)} file uploads failed")
                
                # Retry
                retry_count += 1
                if retry_count < max_retries:
                    wait_time = 2 ** retry_count  # Exponential backoff: 2, 4, 8 seconds
                    print(f"   Retrying in {wait_time} seconds...")
                    await asyncio.sleep(wait_time)
        
        except Exception as e:
            print(f"❌ Background upload error (attempt {retry_count + 1}): {e}")
            import traceback
            traceback.print_exc()
            
            retry_count += 1
            if retry_count < max_retries:
                wait_time = 2 ** retry_count
                print(f"   Retrying in {wait_time} seconds...")
                await asyncio.sleep(wait_time)
    
    # All retries failed - log and give up (silent failure)
    print(f"❌ Background upload failed for chat {chat_id} after {max_retries} attempts")
    print(f"   User can still chat - vectorstore is in memory")

def _language_for_prompt(lang_code: str) -> str:
    """
    Map short language codes to clearer labels for prompt conditioning.
    Keeps default as-is if unknown.
    """
    if not lang_code:
        return "English"
    code = lang_code.lower()
    mapping = {
        "en": "English",
        "english": "English",
        "fr": "French",
        "french": "French",
        "es": "Spanish",
        "spanish": "Spanish"
    }
    return mapping.get(code, lang_code)


async def extract_file_insights_from_text(
    text: str, 
    filename: str, 
    chat_id: str, 
    file_id: str, 
    updates: list,
    language: str = "english"
) -> dict:
    """
    Extract key topics and concepts from document text using random sampling.
    Runs in parallel with embedding for speed.
    
    Args:
        text: Full document text
        filename: Name of the file
        chat_id: Chat ID
        file_id: File ID for progress updates
        updates: List to append progress updates to
        language: Browser language for localized insights
    
    Returns:
        Dict with topics, concepts, and document_type
    """
    try:
        updates.append({
            "type": "insight_extraction_start",
            "file_id": file_id,
            "filename": filename
        })
        
        # Sample random sections for fast analysis
        text_length = len(text)
        
        if text_length < 5000:
            # Small file - use all text
            sample_text = text
        else:
            # Large file - sample 3 random sections (1000 chars each)
            import random
            samples = []
            for _ in range(3):
                start_pos = random.randint(0, max(0, text_length - 1000))
                samples.append(text[start_pos:start_pos + 1000])
            sample_text = "\n\n---\n\n".join(samples)
        
        # Convert to a clearer label for the prompt (e.g., "fr" -> "French")
        prompt_language = _language_for_prompt(language)

        # Use GPT-4o-mini for fast, cheap analysis                
        llm = ChatOpenAI(model="gpt-4o-mini", temperature=0.3)
        
        # Determine response language
        prompt = f"""Analyze this document and identify its CORE PURPOSE in {prompt_language}.

        Document: {filename}
        Content sample:
        {sample_text[:2500]}

        CRITICAL: Focus on the MAIN THESIS or CENTRAL QUESTION of this document.
        - What is the document trying to teach or prove?
        - What is the key relationship or concept being explored?
        - Ignore metadata (demographics, methodology details, sample sizes) - focus on the CONCLUSION or MAIN TEACHING POINT.

        Example of what we want:
        - Document about sleep and testosterone → Topic: "Sleep deprivation reduces testosterone levels"
        - Document about ABCDE assessment → Topic: "Emergency patient assessment protocol"
        - NOT: "Demographics, education levels, sample characteristics" (these are details, not the core topic)

        Identify:
        1. Core topic (1-3 MAIN subjects this document is fundamentally about - the central thesis)
        2. Key concepts (5-10 important terms or findings the student needs to remember)
        3. Document type (research paper, textbook, clinical guide, lecture notes, etc.)
        4. Key insights - For each core topic, extract what the student MUST learn:
           - topic: The main subject (e.g. "Sleep and Testosterone", "ABCDE Assessment")
           - insight: The key finding or teaching point (1 sentence - what should the student remember?)
           - key_points: 2-4 specific facts, steps, or conclusions FROM the document
           - context: Where/when this knowledge applies

        IMPORTANT:
        - Topics should answer "What is this document ABOUT?" not "What variables were measured?"
        - key_points should be actionable knowledge, not methodology details

        Return ONLY valid JSON with content in {prompt_language}:
        {{
        "topics": ["Core topic 1", "Core topic 2"],
        "concepts": ["key term 1", "key term 2", ...],
        "document_type": "type",
        "insights": [
            {{
                "topic": "Main subject of document",
                "insight": "The key finding or teaching point",
                "key_points": ["Important fact 1", "Important fact 2", "Important fact 3"],
                "context": "Where this knowledge is applied"
            }}
        ]
        }}
        """
        
        response = await llm.ainvoke([
            {"role": "system", "content": f"You are an expert tutor who identifies the CORE PURPOSE and MAIN THESIS of educational documents. Focus on what the student needs to LEARN, not on research methodology or metadata. Return only valid JSON with all content in {prompt_language}."},
            {"role": "user", "content": prompt}
        ])
        
        # Parse response
        try:
            insights = json.loads(response.content.strip().strip("```json").strip("```"))
        except json.JSONDecodeError:
            print(f"⚠️ Failed to parse insights JSON for {filename}")
            default_topic = "contenu médical" if language.lower() in ["fr", "french", "français"] else "medical content"
            default_type = "document"
            insights = {
                "topics": [default_topic],
                "concepts": [],
                "document_type": default_type,
                "insights": []
            }
        
        print(f"✅ Extracted insights from {filename}:")
        print(f"   Topics: {insights.get('topics', [])}")
        print(f"   Concepts: {insights.get('concepts', [])[:3]}...")
        print(f"   Educational insights: {len(insights.get('insights', []))} generated")

        # Stream insight batch to frontend
        updates.append({
            "type": "insight_batch",
            "file_id": file_id,
            "filename": filename,
            "topics": insights.get("topics", []),
            "concepts": insights.get("concepts", [])[:5],  # Limit to 5 for UX
            "document_type": insights.get("document_type", ""),
            "insights": insights.get("insights", [])[:3]  # Limit to 3 educational insights
        })
        
        return insights
        
    except Exception as e:
        print(f"⚠️ Insight extraction failed for {filename}: {e}")
        return None

async def firebase_upload_task_simple(file_bytes: bytes, filename: str, chat_id: str):
    """Simple file upload to Firebase Storage (for background task)."""
    try:
        bucket = storage.bucket()
        blob = bucket.blob(f"chats/{chat_id}/uploads/{filename}")
        
        blob.upload_from_string(
            file_bytes,
            content_type=get_content_type(filename)
        )
        
        blob.make_public()
        firebase_url = blob.public_url
        
        print(f"✅ Background: Uploaded {filename} to Firebase Storage")
        return firebase_url
        
    except Exception as e:
        print(f"❌ Background: Failed to upload {filename}: {e}")
        raise

async def process_file_from_bytes(file_bytes: bytes, filename: str, chat_id: str, user_id: str, language: str = "english"):
    """Process a file from bytes and return list of progress updates."""
    updates = []
    file_id = str(uuid4())
    temp_path = None
    documents_for_vectorstore = []
    
    try:
        file_size = len(file_bytes)
        
        updates.append({
            "type": "file_start",
            "file_id": file_id,
            "filename": filename,
            "size": file_size
        })
        
        # Save to temp file
        temp_path = await save_temp_file(file_bytes, filename)
        
        updates.append({
            "type": "file_processing",
            "file_id": file_id,
            "stage": "saved_temp"
        })
        
        # ========================================
        # ONLY EMBEDDING NOW - NO FIREBASE UPLOAD
        # ========================================
        embedding_result = await embed_document_task(
            temp_path, filename, chat_id, file_id, updates, language
        )
        
        # Handle errors
        if isinstance(embedding_result, Exception):
            print(f"❌ Embedding error: {embedding_result}")
            updates.append({
                "type": "embedding_error",
                "file_id": file_id,
                "message": str(embedding_result)
            })
            embedding_result = {"word_count": 0, "chunks": 0, "documents": []}
        
        # Extract documents for vectorstore
        documents_for_vectorstore = embedding_result.get("documents", [])
        
        # 3. Extract insights
        insights = embedding_result.get("insights")
        
        # Final update - embedding complete
        updates.append({
            "type": "file_complete",
            "file_id": file_id,
            "filename": filename,
            "word_count": embedding_result.get("word_count", 0),
            "chunk_count": embedding_result.get("chunks", 0),
            "status": "success"
        })
        
        # Return both updates and documents + file_bytes for background upload
        return {
            "updates": updates,
            "documents": documents_for_vectorstore,
            "filename": filename,
            "file_bytes": file_bytes,  # ← Include for background upload
            "insights": insights 
        }
        
    except Exception as e:
        print(f"❌ Error processing {filename}: {e}")
        import traceback
        traceback.print_exc()
        
        updates.append({
            "type": "file_error",
            "file_id": file_id,
            "filename": filename,
            "message": str(e)
        })
        
        return {
            "updates": updates,
            "documents": [],
            "filename": filename,
            "file_bytes": file_bytes
        }
    
    finally:
        # Clean up temp file
        if temp_path and os.path.exists(temp_path):
            try:
                os.unlink(temp_path)
                print(f"✅ Cleaned up temp file: {temp_path}")
            except Exception as cleanup_error:
                print(f"⚠️ Failed to cleanup {temp_path}: {cleanup_error}")

async def embed_document_task(temp_path: str, filename: str, chat_id: str, file_id: str, updates: list, language: str = "english"):
    """Embed document and extract insights in parallel."""
    try:
        updates.append({
            "type": "embedding_start",
            "file_id": file_id
        })
        
        # Verify file exists
        if not os.path.exists(temp_path):
            raise FileNotFoundError(f"Temp file not found: {temp_path}")
        
        print(f"📄 Loading document: {filename}")
        
        # Load document
        loader = get_loader_for_file(temp_path)
        pages = loader.load()
        
        print(f"✅ Loaded {len(pages)} pages from {filename}")
        
        updates.append({
            "type": "embedding_progress",
            "file_id": file_id,
            "stage": "loaded_pages",
            "page_count": len(pages)
        })
        
        # Extract text
        text = "\n\n".join([p.page_content for p in pages])
        word_count = len(text.split())
        
        print(f"📝 Extracted {word_count} words from {filename}")
        
        # ========================================
        # 🆕 START INSIGHT EXTRACTION IN PARALLEL
        # ========================================
        insight_task = asyncio.create_task(
            extract_file_insights_from_text(text, filename, chat_id, file_id, updates, language)
        )
        
        # ========================================
        # CONTINUE WITH EMBEDDING (PARALLEL)
        # ========================================
        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=1000,
            chunk_overlap=200
        )
        chunks = text_splitter.split_text(text)
        
        print(f"✂️ Split into {len(chunks)} chunks")
        
        updates.append({
            "type": "embedding_progress",
            "file_id": file_id,
            "stage": "chunked",
            "chunk_count": len(chunks)
        })
        
        # Create documents
        documents = [
            Document(page_content=chunk, metadata={"source": filename})
            for chunk in chunks
        ]
        
        print(f"🔤 Creating embeddings for {len(documents)} documents...")
        
        # Get embeddings
        embeddings = OpenAIEmbeddings()
        
        # Get session
        if chat_id not in ACTIVE_SESSIONS:
            print(f"⚠️ No session found for {chat_id}, creating...")
            ACTIVE_SESSIONS[chat_id] = NursingTutor(chat_id)

        session = ACTIVE_SESSIONS[chat_id]
        print(f"📤 Upload using session object id: {id(session.session)}")

        # Add to combined vectorstore
        if session.session.vectorstore:
            print(f"➕ Adding to existing vectorstore")
            session.session.vectorstore.add_documents(documents)
        else:
            print(f"🆕 Creating new vectorstore")
            session.session.vectorstore = FAISS.from_documents(documents, embeddings)

        print(f"✅ Embedding complete for {filename}")
        print(f"📤 Vectorstore now set: {session.session.vectorstore is not None}")
        
        updates.append({
            "type": "embedding_complete",
            "file_id": file_id,
            "word_count": word_count,
            "chunks": len(documents)
        })
        
        # ========================================
        # 🆕 WAIT FOR INSIGHTS (SHOULD BE READY)
        # ========================================
        try:
            insights = await asyncio.wait_for(insight_task, timeout=20.0)
        except asyncio.TimeoutError:
            print(f"⚠️ Insight extraction timed out for {filename}")
            insights = None
        
        # Return documents + insights
        return {
            "word_count": word_count,
            "chunks": len(documents),
            "documents": documents,
            "insights": insights  # ← Include insights
        }
        
    except Exception as e:
        print(f"❌ Embedding error for {filename}: {e}")
        import traceback
        traceback.print_exc()
        raise

async def firebase_upload_task(file_bytes, filename, chat_id, file_id, updates):
    """Upload to Firebase - appends progress to updates list"""
    try:
        updates.append({
            "type": "firebase_start",
            "file_id": file_id
        })
        
        bucket = storage.bucket()
        blob = bucket.blob(f"chats/{chat_id}/uploads/{filename}")
        
        # Upload
        blob.upload_from_string(
            file_bytes,
            content_type=get_content_type(filename)
        )
        
        # Make public and get URL
        blob.make_public()
        firebase_url = blob.public_url
        
        updates.append({
            "type": "firebase_complete",
            "file_id": file_id,
            "firebase_url": firebase_url
        })
        
        return firebase_url
        
    except Exception as e:
        print(f"Firebase upload error for {filename}: {e}")
        raise
         
def get_content_type(filename: str) -> str:
    """
    Get MIME type from filename using Python's standard library.
    
    Args:
        filename: Name of the file (e.g., "notes.pdf")
    
    Returns:
        MIME type string (e.g., "application/pdf")
    """
    # Guess MIME type from filename
    content_type, _ = mimetypes.guess_type(filename)
    
    # If unknown, default to generic binary
    return content_type or 'application/octet-stream'

async def save_temp_file(file_bytes: bytes, filename: str) -> str:
    """
    Save uploaded file bytes to a temporary file.
    
    Args:
        file_bytes: The file content as bytes
        filename: Original filename (used to get extension)
    
    Returns:
        Path to the temporary file
    """
    temp_dir = get_temp_dir()
    suffix = os.path.splitext(filename)[-1]  # Get file extension (.pdf, .docx, etc.)
    
    # Create a temporary file with the correct extension
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix, dir=temp_dir) as f:
        f.write(file_bytes)
        temp_path = f.name
    
    print(f"✅ Saved temp file: {temp_path}")
    return temp_path

async def ensure_session_with_vectorstore(chat_id: str):
    """
    Ensure session exists and has vectorstore loaded.
    Shows "reading documents" message when downloading from Firebase.
    """
    
    if chat_id in ACTIVE_SESSIONS:
        print(f"✅ Session already exists for {chat_id}")
        return ACTIVE_SESSIONS[chat_id]
    
    print(f"🆕 Creating new session for {chat_id}")
    
    # Create session
    ACTIVE_SESSIONS[chat_id] = NursingTutor(chat_id)
    
    # Try to load existing vectorstore
    print(f"📚 Checking for existing documents...")
    vectorstore = await vectorstore_manager.load_combined_vectorstore_from_firebase(chat_id)
    
    if vectorstore:
        ACTIVE_SESSIONS[chat_id].session.vectorstore = vectorstore
        print(f"✅ Loaded existing vectorstore into session")
    else:
        print(f"📝 No existing vectorstore, will create new one")
    
    return ACTIVE_SESSIONS[chat_id]

# ============================================================================
# HELPERS FOR FILE UPLOAD END
# ============================================================================

@app.post("/chat/generate-summary")
async def generate_summary(request:SummaryRequest):
    
    from tools.quiztools import _search_vectorstore_for_summary,set_session_context
     # Get or create session for this chat
    if request.chat_id not in ACTIVE_SESSIONS:
        ACTIVE_SESSIONS[request.chat_id] = NursingTutor(request.chat_id)
        print(f"Created new session for chat_id: {request.chat_id}")
    
    # GET TUTOR FOR CURRENT SESSION
    nursing_tutor = ACTIVE_SESSIONS[request.chat_id]
    print(f"nursing tutor created")
        
    
    nursing_tutor.session.name_last_document_used=request.filename
    
    # === FIX FOR ATTRIBUTE ERROR ===
    # The NursingTutor object seems to be missing the chat_id attribute internally.
    # We enforce its existence here before passing it to set_session_context.
    if not hasattr(nursing_tutor, 'chat_id'):
        nursing_tutor.chat_id = request.chat_id
        
    set_session_context(nursing_tutor)
    
    from tools.quiztools import _search_vectorstore_for_summary
    chunks = await _search_vectorstore_for_summary(request.filename, request.chat_id, "", "detailed")
    
    print("Got the chunks for summary through endpoint")
    
    async def json_chunk_generator():
       async for chunk in   nursing_tutor.stream_document_summary(
           relevant_chunks=chunks,
           detail_level="detailed",
           filename=request.filename,
           language=request.language):
            yield json.dumps({
                "answer_chunk": chunk
            }) + "\n"
    
    print("Streaming summary through endpoint response")
    
    return StreamingResponse(
        json_chunk_generator(),
        media_type="application/json"
    )
 
@app.post("/plan")
async def create_plan(request: PlanRequest):
    from tools.quiztools import search_documents
     
    print("building plan for study guide",request)
    
    session = ACTIVE_SESSIONS[request.chat_id]
    
    try:
        # Get context using your existing tool
        search_result = await search_documents.ainvoke({
            "query": request.topic
        })
        
        context = search_result.get("context", "")
        
        
        # STEP 1: Create a prompt asking LLM to generate a plan
        prompt = f"""
        Create {request.num_sections} sections for a study guide about {request.topic}.
        base in this context {context}
        in this language {session.session.user_language}
        Return ONLY a JSON array:
        [
        {{"id": "introduction", "title": "Introduction", "color": "blue"}},
        {{"id": "concepts", "title": "Key Concepts", "color": "green"}}
        ]
        """
        # This creates an LLM instance
        llm = ChatOpenAI(
            model="gpt-4o",           # Larger context for richer plans
            temperature=0.3
        )
        
        # STEP 2: Send prompt to LLM (this generates the actual plan)
        response = await llm.ainvoke([{"role": "user", "content": prompt}])

        # STEP 3: LLM returns a string that looks like JSON
        plan_json = response.content.strip()
        # plan_json is now a STRING: '[{"id":"intro","title":"Introduction"...}]'

        # STEP 4: Clean up markdown code blocks if LLM wrapped it
        if plan_json.startswith("```"):
            plan_json = plan_json.split("```")[1]
            if plan_json.startswith("json"):
                plan_json = plan_json[4:]
            plan_json = plan_json.strip()

        # STEP 5: Convert JSON string to Python list/dict
        sections = json.loads(plan_json)  # ← This parses the string into actual Python objects
        # sections is now a Python LIST: [{"id": "intro", "title": "Introduction"...}]
        # Generate plan...
        
        # Return BOTH sections AND context
        plan = {
            "sections": sections,
            "context": context  # ← Add this
        }
    
        print("This is the plan", plan)
        
        return plan
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    
@app.post("/search")
async def search_for_context(request: dict):
    """
    Reuse your existing search_documents tool to get context
    """
    from tools.quiztools import search_documents
    
    try:
        # Call your existing search tool
        result = await search_documents.ainvoke({
            "query": request.get("query"),
            # Add any other params your search tool needs
        })
        
        # Extract context string
        context = result.get("context", "")
        
        return {"context": context}
        
    except Exception as e:
        print(f"Search error: {e}")
        raise HTTPException(status_code=500, detail=str(e))
 
@app.post("/generate-section")
async def generate_section(request: SectionRequest):
    """Generate content for a section using RAG context"""
    
    print("GENERATING SECTION BASED ON",request)
    
    try:
        llm = ChatOpenAI(model="gpt-4o", temperature=0.3)    
        prompt = f"""
        You are creating educational content for a study guide.

        Topic: {request.topic}
        Section: {request.section_title}

        Retrieved Context from Student's Documents:
        {request.context}

        Generate comprehensive educational content for this section using ONLY information from the documents.

        Format as HTML:
        - <h3>Subsection Title</h3>
        - <p>Explanation text</p>
        - <ul><li>Bullet points</li></ul>

        - <div class="card card-blue">
            <div class="card-title">🔑 Key Concept</div>
            <p>Important nursing information</p>
        </div>

        - <div class="card card-green">
            <div class="card-title">✅ Clinical Application</div>
            <p>How to apply this in practice</p>
        </div>

        - <div class="card card-yellow">
            <div class="card-title">⚠️ Critical Alert</div>
            <p>Warning or safety information</p>
        </div>

        - Use <strong> for emphasis
        - Use <span class="highlight">term</span> for key terms

        Guidelines:
        - Use nursing emojis (🩺 💊 🫁 ❤️ 🧠)
        - Include rationales (WHY, not just WHAT)
        - Focus on NCLEX-style critical thinking
        - Be comprehensive but concise
        - Use proper medical terminology
        
        IMPORTANT:
        Make sure all the content written are in the same language, either french or english
        It should be uniform from the header titles, to critical alert, key concepts etc

        Return ONLY the HTML content, no markdown code blocks.
        """
        
        response = await llm.ainvoke([{"role": "user", "content": prompt}])
        content = response.content.strip()
        
        # Clean up
        if content.startswith("```"):
            content = content.split("```")[1]
            if content.startswith("html"):
                content = content[4:]
            content = content.strip()
        
        return {"content": content}
        
    except Exception as e:
        print(f"Error generating section: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ============================================================================
# STUDY MODE ENDPOINTS
# Duolingo-style learning path generation
# ============================================================================

from models.requests import StudyPlanRequest, StudyItemRequest
import hashlib

@app.post("/study/plan")
async def generate_study_plan(request: StudyPlanRequest):
    """
    Generate a personalized study path based on uploaded documents.

    HOW IT WORKS:
    1. Get file insights (topics, concepts) from the session
    2. Use LLM to create a logical learning sequence
    3. Return a list of nodes (lesson, flashcard, quiz, audio)

    Each node has:
    - id: unique identifier
    - type: "lesson" | "flashcard" | "quiz" | "audio"
    - label: topic name shown to user
    - tags: context tags for content generation
    - difficulty: 1-3 scale
    - status: "locked" | "available" | "completed"

    Frontend will store this in Firestore and track progress.
    """
    print(f"\n{'='*60}")
    print(f"📚 STUDY PATH GENERATION - chat_id: {request.chat_id}")
    print(f"{'='*60}")

    try:
        # ------------------------------------------
        # STEP 1: Get or create session & load insights
        # ------------------------------------------
        if request.chat_id not in ACTIVE_SESSIONS:
            ACTIVE_SESSIONS[request.chat_id] = NursingTutor(request.chat_id)
            await ACTIVE_SESSIONS[request.chat_id].load_file_insights_from_firebase()
            print(f"🆕 Created new session for study plan")

        session = ACTIVE_SESSIONS[request.chat_id]
        file_insights = getattr(session.session, "file_insights", {})

        # ------------------------------------------
        # STEP 2: Collect all topics and concepts
        # ------------------------------------------
        all_topics = []
        all_concepts = []

        for filename, insights in file_insights.items():
            if insights:
                all_topics.extend(insights.get("topics", []))
                all_concepts.extend(insights.get("concepts", []))

        # Remove duplicates
        unique_topics = list(set(all_topics))[:8]  # Max 8 topics for manageable path
        unique_concepts = list(set(all_concepts))[:15]

        print(f"📊 Found {len(unique_topics)} topics: {unique_topics}")
        print(f"📊 Found {len(unique_concepts)} concepts")

        # ------------------------------------------
        # STEP 3: Fallback if no insights found
        # ------------------------------------------
        if not unique_topics:
            # Try to get context from vectorstore
            if session.session.vectorstore:
                docs = session.session.vectorstore.similarity_search("main topics", k=5)
                context_text = "\n".join([doc.page_content[:500] for doc in docs])

                # Quick LLM call to extract topics
                llm = ChatOpenAI(model="gpt-4o-mini", temperature=0.3)
                topic_prompt = f"""Extract 3-5 main study topics from this content. Return as JSON array of strings.

Content:
{context_text[:2000]}

Return ONLY: ["Topic 1", "Topic 2", "Topic 3"]"""

                topic_response = await llm.ainvoke([{"role": "user", "content": topic_prompt}])
                try:
                    unique_topics = json.loads(topic_response.content.strip())
                except:
                    unique_topics = ["General Medical Knowledge"]

        # ------------------------------------------
        # STEP 4: Generate learning path with LLM
        # ------------------------------------------
        prompt_language = _language_for_prompt(request.language)

        llm = ChatOpenAI(model="gpt-4o-mini", temperature=0.5)

        path_prompt = f"""You are an expert learning path designer. Create a Duolingo-style study path for these topics.

TOPICS TO COVER: {unique_topics}
KEY CONCEPTS: {unique_concepts[:10]}

RULES:
1. Create 8-12 nodes total
2. Start with a LESSON to introduce the first topic
3. Follow each lesson with FLASHCARDS or QUIZ to reinforce
4. Include at least one AUDIO node for variety
5. Progress from easier to harder content
6. End with a comprehensive QUIZ

NODE TYPES:
- "lesson": Short explanation (5-8 lines) with key points
- "flashcard": Single front/back card for memorization
- "quiz": Multiple choice question with rationale
- "audio": Topic for audio explanation (1-2 min)

Return ONLY valid JSON array in {prompt_language}:
[
  {{"id": "node_1", "type": "lesson", "label": "Introduction to [Topic]", "tags": ["intro", "basics"], "difficulty": 1}},
  {{"id": "node_2", "type": "flashcard", "label": "Key Term: [Term]", "tags": ["vocabulary"], "difficulty": 1}},
  {{"id": "node_3", "type": "quiz", "label": "Quick Check: [Topic]", "tags": ["assessment"], "difficulty": 1}},
  ...
]

IMPORTANT:
- "label" should be clear and specific (what the user will see)
- "tags" help with content generation later
- "difficulty" is 1 (easy), 2 (medium), or 3 (hard)
- Labels should be in {prompt_language}"""

        response = await llm.ainvoke([{"role": "user", "content": path_prompt}])

        # Parse the response
        path_json = response.content.strip()

        # Clean markdown code blocks if present
        if path_json.startswith("```"):
            path_json = path_json.split("```")[1]
            if path_json.startswith("json"):
                path_json = path_json[4:]
            path_json = path_json.strip()

        nodes = json.loads(path_json)

        # ------------------------------------------
        # STEP 5: Add status and validate structure
        # ------------------------------------------
        for i, node in enumerate(nodes):
            # First node is available, rest are locked
            node["status"] = "available" if i == 0 else "locked"

            # Ensure required fields exist
            if "id" not in node:
                node["id"] = f"node_{i+1}"
            if "difficulty" not in node:
                node["difficulty"] = 1 + (i // 4)  # Gradually increase
            if "tags" not in node:
                node["tags"] = []

        print(f"✅ Generated study path with {len(nodes)} nodes")
        for node in nodes:
            print(f"   - {node['type']}: {node['label']}")

        # ------------------------------------------
        # STEP 6: Return the study path
        # ------------------------------------------
        return {
            "nodes": nodes,
            "topics": unique_topics,
            "total_nodes": len(nodes),
            "estimated_time_minutes": len(nodes) * 3  # ~3 min per node
        }

    except Exception as e:
        print(f"❌ Study path generation failed: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/study/generate-item")
async def generate_study_item(request: StudyItemRequest):
    """
    Generate content for a single study node on-demand.

    WHY ON-DEMAND:
    - Saves cost (only generate what user actually views)
    - Feels more dynamic and personalized
    - Uses anti-repeat hashes to avoid duplicate content

    CONTENT TYPES:
    - lesson: {title, body, keyPoints[]}
    - flashcard: {front, back}
    - quiz: {question, options[], correctIndex, rationale}
    - audio: {topic, intent, suggestedDuration}

    Returns:
    - type: the node type
    - content: the generated content object
    - hash: content hash for anti-repeat tracking
    """
    print(f"\n{'='*60}")
    print(f"📝 STUDY ITEM GENERATION - {request.node_type}: {request.node_label}")
    print(f"{'='*60}")

    try:
        # ------------------------------------------
        # STEP 1: Get session and vectorstore context
        # ------------------------------------------
        if request.chat_id not in ACTIVE_SESSIONS:
            ACTIVE_SESSIONS[request.chat_id] = NursingTutor(request.chat_id)
            await ACTIVE_SESSIONS[request.chat_id].load_file_insights_from_firebase()

        session = ACTIVE_SESSIONS[request.chat_id]

        # Load vectorstore if needed
        if session.session.vectorstore is None:
            loaded_vs = await vectorstore_manager.load_combined_vectorstore_from_firebase(request.chat_id)
            if loaded_vs:
                session.session.vectorstore = loaded_vs

        # ------------------------------------------
        # STEP 2: Search for relevant context
        # ------------------------------------------
        context = ""
        if session.session.vectorstore:
            # Search for content related to this node's topic
            search_query = f"{request.node_label} {' '.join(request.context_tags)}"
            docs = session.session.vectorstore.similarity_search(search_query, k=4)
            context = "\n\n".join([doc.page_content for doc in docs])
            print(f"📚 Found {len(docs)} relevant chunks for context")

        # ------------------------------------------
        # STEP 3: Generate content based on type
        # ------------------------------------------
        prompt_language = _language_for_prompt(request.language)
        llm = ChatOpenAI(model="gpt-4o-mini", temperature=0.6)

        content = None

        if request.node_type == "lesson":
            content = await _generate_lesson(
                llm, request.node_label, context, prompt_language, request.asked_hashes
            )

        elif request.node_type == "flashcard":
            content = await _generate_flashcard(
                llm, request.node_label, context, prompt_language, request.asked_hashes
            )

        elif request.node_type == "quiz":
            content = await _generate_quiz_item(
                llm, request.node_label, context, prompt_language, request.asked_hashes
            )

        elif request.node_type == "audio":
            content = await _generate_audio_config(
                llm, request.node_label, context, prompt_language
            )

        else:
            raise HTTPException(status_code=400, detail=f"Unknown node type: {request.node_type}")

        # ------------------------------------------
        # STEP 4: Generate content hash for anti-repeat
        # ------------------------------------------
        content_str = json.dumps(content, sort_keys=True)
        content_hash = hashlib.md5(content_str.encode()).hexdigest()[:12]

        print(f"✅ Generated {request.node_type} content (hash: {content_hash})")

        return {
            "type": request.node_type,
            "content": content,
            "hash": content_hash
        }

    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Study item generation failed: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


# ============================================================================
# STUDY ITEM GENERATION HELPERS
# Each function generates a specific content type
# ============================================================================

async def _generate_lesson(llm, label: str, context: str, language: str, asked_hashes: list) -> dict:
    """
    Generate a short lesson (5-8 lines) with key takeaways.

    Think of it like a Duolingo "tip" before exercises.
    """
    prompt = f"""Create a SHORT lesson about: {label}

CONTEXT FROM STUDENT'S DOCUMENTS:
{context[:2000] if context else "No specific context - use general knowledge"}

REQUIREMENTS:
1. Title: Clear, engaging title
2. Body: 5-8 sentences explaining the topic
3. Key Points: 3 bullet points to remember

Write in {language}. Be concise but educational.
Focus on the "why" not just "what".

Return ONLY valid JSON:
{{
  "title": "Lesson title here",
  "body": "The explanation paragraph here. Keep it clear and engaging...",
  "keyPoints": ["Point 1", "Point 2", "Point 3"]
}}"""

    response = await llm.ainvoke([{"role": "user", "content": prompt}])

    # Parse response
    content_json = response.content.strip()
    if content_json.startswith("```"):
        content_json = content_json.split("```")[1]
        if content_json.startswith("json"):
            content_json = content_json[4:]
        content_json = content_json.strip()

    return json.loads(content_json)


async def _generate_flashcard(llm, label: str, context: str, language: str, asked_hashes: list) -> dict:
    """
    Generate 5 high-quality flashcards for a single node using the enhanced generator.

    Uses the _generate_single_flashcard from flashcard_tools which produces:
    - Scannable answers with bullets and bold formatting
    - Topic assignment for organization
    - Better deduplication across cards

    Returns an object with a 'cards' array containing 5 flashcards.
    Each card has front/back/topic.
    """
    from tools.flashcard_tools import _generate_single_flashcard

    # Build content context for the generator
    if context:
        content_context = f"Document content:\n{context[:15000]}"
    else:
        content_context = f"""You are generating flashcards about: {label}

        Cover key concepts, definitions, and important facts that students need to memorize.
        If this is a broad topic, ensure diverse coverage of subtopics."""

    cards = []
    generated_fronts = []
    existing_topics = []

    # Generate 5 flashcards one at a time for better quality
    for card_num in range(1, 6):
        flashcard_data = await _generate_single_flashcard(
            content=content_context,
            topic=label,
            card_num=card_num,
            language=language,
            cards_to_avoid=generated_fronts,
            existing_topics=existing_topics
        )

        if flashcard_data:
            # Track front for deduplication
            generated_fronts.append(flashcard_data['front'])

            # Track topic for consistency
            card_topic = flashcard_data.get('topic', '')
            if card_topic and card_topic not in existing_topics:
                existing_topics.append(card_topic)

            cards.append(flashcard_data)
            print(f"✅ Flashcard {card_num}/5 generated - Topic: {flashcard_data.get('topic', 'N/A')}")

    # Ensure we have at least some cards
    if not cards:
        # Fallback to simple generation if enhanced method fails
        print("⚠️ Enhanced flashcard generation failed, using fallback")
        cards = [
            {"front": f"What is {label}?", "back": f"Definition of {label}...", "topic": label},
            {"front": f"Key features of {label}?", "back": f"The key features are...", "topic": label},
            {"front": f"Why is {label} important?", "back": f"It is important because...", "topic": label},
            {"front": f"How is {label} applied?", "back": f"It is applied by...", "topic": label},
            {"front": f"Common mistakes with {label}?", "back": f"Common mistakes include...", "topic": label}
        ]

    print(f"📇 Generated {len(cards)} high-quality flashcards")
    return {"cards": cards}


async def _generate_quiz_item(llm, label: str, context: str, language: str, asked_hashes: list) -> dict:
    """
    Generate 5 high-quality multiple choice questions for a single node.

    Uses the _generate_single_question from quiztools which produces:
    - Random answer positioning (not always B or C)
    - Topic assignment for organization
    - Better justifications with bold formatting
    - Knowledge mode (factual recall, not NCLEX scenarios)

    Returns an object with a 'questions' array containing 5 quiz questions.
    Each question has: question, options, correctIndex, rationale, topic.
    """
    import random
    from tools.quiztools import _generate_single_question

    # Build content context for the generator
    if context:
        content_context = f"Document content:\n{context[:12000]}"
    else:
        content_context = f"""You are generating questions about: {label}

            If this is a broad topic (like 'research design', 'pharmacology', 'cardiac care'),
            ensure you test diverse subtopics and concepts within that domain."""

    questions = []
    generated_question_texts = []
    existing_topics = []

    # Generate 5 questions one at a time for better quality
    for question_num in range(1, 6):
        # Random answer position for each question
        random_target_letter = random.choice(['A', 'B', 'C', 'D'])

        question_data = await _generate_single_question(
            content=content_context,
            topic=label,
            difficulty="medium",
            question_num=question_num,
            language=language,
            questions_to_avoid=generated_question_texts,
            target_letter=random_target_letter,
            existing_topics=existing_topics,
            quiz_mode="knowledge"  # Use knowledge mode for study path (factual, not NCLEX scenarios)
        )

        if question_data:
            # Track question text for deduplication
            generated_question_texts.append(question_data['question'])

            # Track topic for consistency
            question_topic = question_data.get('topic', '')
            if question_topic and question_topic not in existing_topics:
                existing_topics.append(question_topic)

            # Convert format for frontend compatibility
            # Frontend expects: { question, options, correctIndex, rationale }
            answer = question_data.get('answer', 'A)')
            answer_letter = answer[0] if answer else 'A'
            correct_index = ord(answer_letter) - ord('A')

            formatted_question = {
                "question": question_data['question'],
                "options": question_data['options'],
                "correctIndex": correct_index,
                "rationale": question_data.get('justification', ''),
                "topic": question_data.get('topic', label)
            }

            questions.append(formatted_question)
            print(f"✅ Question {question_num}/5 generated - Answer: {answer_letter}, Topic: {formatted_question['topic']}")

    # Ensure we have at least some questions
    if not questions:
        # Fallback to simple generation if enhanced method fails
        print("⚠️ Enhanced quiz generation failed, using fallback")
        questions = [
            {
                "question": f"What is the primary purpose of {label}?",
                "options": ["A) Option 1", "B) Option 2", "C) Option 3", "D) Option 4"],
                "correctIndex": 0,
                "rationale": f"Option A is correct because it describes {label}.",
                "topic": label
            }
        ]

    print(f"❓ Generated {len(questions)} high-quality quiz questions")
    return {"questions": questions}


async def _generate_audio_config(llm, label: str, context: str, language: str) -> dict:
    """
    Generate configuration for an audio lesson.

    This returns the TOPIC and INTENT, not the actual audio.
    The frontend will use this to call the existing audio generation endpoint.
    """
    prompt = f"""Create an audio lesson config for: {label}

CONTEXT:
{context[:1000] if context else "No specific context"}

Return a configuration for a 1-2 minute audio explanation.

Return ONLY valid JSON:
{{
  "topic": "Clear topic title for audio",
  "intent": "teach",
  "suggestedDuration": 2
}}"""

    response = await llm.ainvoke([{"role": "user", "content": prompt}])

    content_json = response.content.strip()
    if content_json.startswith("```"):
        content_json = content_json.split("```")[1]
        if content_json.startswith("json"):
            content_json = content_json[4:]
        content_json = content_json.strip()

    return json.loads(content_json)


@app.post("/chat/generate-title", response_model=GenerateTitleResponse)
async def generate_chat_title(request: GenerateTitleRequest):
    
    if request.message:
        prompt = PromptTemplate(
            template="""
            You are an AI assistant for nursing students.

            Based on the following user message, generate a short, descriptive chat title in 3 to 6 words.

            Requirements:
            - Be concise and clear
            - Write a title in the language of Message you received
            - Max 6 words

            Message:
            {message}
            """,
            input_variables=["message"]
        )
        
       
        # Use gpt-4.1-nano for title generation (simple classification task)
        llm = ChatOpenAI(
            temperature=0.8,
            model="gpt-4.1-nano",
            streaming=False
        )

        chain = prompt | llm | StrOutputParser()
        generated_title = await chain.ainvoke({"message": request.message})
        
        # Clean up the title
        cleaned_title = generated_title.replace('"', '').replace("'", "").strip()

        return GenerateTitleResponse(title=cleaned_title)

# ============================================================================
# QUESTION BANK IMPORT ENDPOINT
# ============================================================================
from pydantic import BaseModel
from typing import List, Optional, Dict, Any

class QuestionMetadata(BaseModel):
    """Metadata object matching LLM output format."""
    sourceLanguage: str = "en"
    topic: str = ""
    category: str = "nursing"
    difficulty: str = "medium"
    correctAnswerIndex: int = 0
    sourceDocument: str = "admin_import"
    keywords: List[str] = []

class QuestionImport(BaseModel):
    """
    Single question for import.
    Matches the EXACT format your LLM generates.
    """
    question: str
    options: List[str]  # ["A) ...", "B) ...", "C) ...", "D) ..."]
    answer: str         # The correct option text (e.g., "B) ...")
    justification: str  # HTML explanation
    topic: str          # e.g., "cardiac medications"
    metadata: Optional[QuestionMetadata] = None  # Optional - matches LLM output

class BulkImportRequest(BaseModel):
    """Request body for bulk question import."""
    questions: List[QuestionImport]

@app.post("/admin/import-questions")
async def import_questions(request: BulkImportRequest):
    """
    Import AI-generated questions into the Question Bank.

    Accepts questions in the EXACT format your LLM generates.
    The metadata object is optional but will be used if provided.

    Example request body (matches LLM output):
    {
        "questions": [
            {
                "question": "A patient receiving digoxin reports nausea...",
                "options": [
                    "A) Continue the medication as prescribed",
                    "B) Hold the medication and notify the provider",
                    "C) Administer an antiemetic",
                    "D) Document the findings and reassess later"
                ],
                "answer": "B) Hold the medication and notify the provider",
                "justification": "<strong>Option B is correct</strong> because...",
                "topic": "Cardiac Medications",
                "metadata": {
                    "sourceLanguage": "en",
                    "topic": "Cardiac Medications",
                    "category": "nursing",
                    "difficulty": "medium",
                    "correctAnswerIndex": 1,
                    "sourceDocument": "conversational_generation",
                    "keywords": ["digoxin", "toxicity", "cardiac"]
                }
            }
        ]
    }
    """
    from services.question_bank import question_bank

    results = {
        "total": len(request.questions),
        "imported": 0,
        "duplicates": 0,
        "errors": [],
        "imported_ids": []
    }

    for i, q in enumerate(request.questions):
        try:
            # Extract language and difficulty from metadata if available
            language = "en"
            difficulty = "medium"

            if q.metadata:
                language = q.metadata.sourceLanguage or "en"
                difficulty = q.metadata.difficulty or "medium"

            # Normalize language code
            if language.lower().startswith("fr"):
                language = "fr"
            elif language.lower().startswith("es"):
                language = "es"
            else:
                language = "en"

            # The question_data format expected by save_question
            question_data = {
                "question": q.question,
                "options": q.options,
                "answer": q.answer,
                "justification": q.justification,
                "topic": q.topic
            }

            # Save to question bank
            doc_id = await question_bank.save_question(
                question_data=question_data,
                topic=q.topic,
                language=language,
                difficulty=difficulty,
                chat_id="admin_import"  # Mark as admin import
            )

            if doc_id:
                results["imported"] += 1
                results["imported_ids"].append(doc_id)
                print(f"✅ Imported question {i+1}: {doc_id}")
            else:
                results["duplicates"] += 1
                print(f"⚠️ Question {i+1} skipped (duplicate)")

        except Exception as e:
            error_msg = f"Question {i+1}: {str(e)}"
            results["errors"].append(error_msg)
            print(f"❌ Error importing question {i+1}: {e}")

    print(f"📊 Import complete: {results['imported']} imported, {results['duplicates']} duplicates, {len(results['errors'])} errors")

    return results


@app.get("/admin/question-bank-stats")
async def get_question_bank_stats():
    """
    Get statistics about the Question Bank.

    Returns total questions, breakdown by language and category.
    """
    from services.question_bank import question_bank

    stats = await question_bank.get_bank_stats()
    return stats


# ============================================================================
# SPEECH-TO-TEXT (Voice Input)
# ============================================================================
from openai import OpenAI

@app.post("/speech-to-text")
async def speech_to_text(audio: UploadFile = File(...)):
    """
    Transcribe audio to text using OpenAI Whisper API.

    Accepts audio files (webm, mp3, wav, m4a, etc.)
    Returns the transcribed text that can be used as a chat prompt.

    Cost: ~$0.006 per minute of audio
    """
    try:
        # Validate file type
        content_type = audio.content_type or ""
        if not any(t in content_type for t in ["audio", "video/webm"]):
            # Also check file extension as fallback
            ext = audio.filename.split(".")[-1].lower() if audio.filename else ""
            if ext not in ["webm", "mp3", "wav", "m4a", "ogg", "mp4"]:
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported audio format: {content_type}. Use webm, mp3, wav, m4a, or ogg."
                )

        # Read the audio file
        audio_bytes = await audio.read()

        if len(audio_bytes) == 0:
            raise HTTPException(status_code=400, detail="Empty audio file")

        # Max 25MB (Whisper API limit)
        if len(audio_bytes) > 25 * 1024 * 1024:
            raise HTTPException(status_code=400, detail="Audio file too large (max 25MB)")

        # Create a temp file for the audio
        ext = audio.filename.split(".")[-1] if audio.filename else "webm"
        with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
            tmp.write(audio_bytes)
            tmp_path = tmp.name

        try:
            # Call OpenAI Whisper API
            client = OpenAI()

            with open(tmp_path, "rb") as audio_file:
                transcript = client.audio.transcriptions.create(
                    model="whisper-1",
                    file=audio_file,
                    response_format="text"
                )

            # Clean up the transcription
            transcribed_text = transcript.strip() if isinstance(transcript, str) else transcript

            print(f"🎤 [STT] Transcribed {len(audio_bytes)} bytes -> '{transcribed_text[:100]}...'")

            return {
                "success": True,
                "text": transcribed_text,
                "audio_size_bytes": len(audio_bytes)
            }

        finally:
            # Clean up temp file
            if os.path.exists(tmp_path):
                os.remove(tmp_path)

    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ [STT] Error: {e}")
        raise HTTPException(status_code=500, detail=f"Transcription failed: {str(e)}")


# ============================================================================
# HEALTH CHECK
# ============================================================================
@app.post("/warm_up")
async def warm_up():
    return {"status": "ok", "message": "Server warmed up successfully"}

# ============================================================================
# RUN THE APP
# ============================================================================
if __name__ == "__main__":
    import uvicorn
    import os
    
    # Get port from environment (Cloud Run sets this to 8080)
    port = int(os.getenv("PORT", 8080))
    
    print(f"Starting server on port {port}...")
    
    uvicorn.run(
        "main:app",
        host="0.0.0.0",  # CRITICAL: Must be 0.0.0.0, not 127.0.0.1
        port=port,
        reload=False,  # CRITICAL: No reload in production
    )
